from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

class ppt:
    def __init__(self, filename="status_report.pptx"):
        self.filename = filename
        self.prs = Presentation()
        self.current_slide = None

    def add_slide(self):
        self.current_slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Blank slide

    def add_title(self, text):
        title_shape = self.current_slide.shapes.title
        if not title_shape:
            title_shape = self.current_slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_shape.text = text
        p = title_shape.text_frame.paragraphs[0]
        run = p.runs[0] if p.runs else p.add_run()
        run.font.size = Pt(28)
        p.alignment = PP_ALIGN.CENTER

    def add_table(self, data):
        if not data:
            return

        if not self.current_slide:
            self.add_slide()

        headers = list(data[0].keys())
        rows, cols = len(data) + 1, len(headers)

        left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 + 0.3 * len(data))
        table = self.current_slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Header row
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(18)

        # Data rows
        for row_idx, row_data in enumerate(data, start=1):
            for col_idx, key in enumerate(headers):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(row_data[key])
                p = cell.text_frame.paragraphs[0]
                # Align second column center, others left
                p.alignment = PP_ALIGN.CENTER if col_idx == 1 else PP_ALIGN.LEFT
                p.font.size = Pt(16)

    def add_graph(self, data, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED):
        if not data:
            return

        if not self.current_slide:
            self.add_slide()

        headers = list(data[0].keys())
        if len(headers) < 2:
            raise ValueError("Data must contain at least two keys per dictionary.")

        label_key, value_key = headers[0], headers[1]
        try:
            categories = [d[label_key] for d in data]
            values = [d[value_key] for d in data]
        except KeyError as e:
            raise ValueError(f"Key '{e.args[0]}' not found in data items.")

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series(value_key, values)

        x, y, cx, cy = Inches(1), Inches(2), Inches(6), Inches(4)
        self.current_slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)

    def save(self):
        self.prs.save(self.filename)

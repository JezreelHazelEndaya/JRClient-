from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR,  MSO_VERTICAL_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

class ppt:
    MAX_CONTENT_HEIGHT = Inches(7.0)      # usable vertical space
    DEFAULT_TOP_OFFSET = Inches(0.5)      # starting top margin
    ELEMENT_SPACING = Inches(0.2)         # space between elements

    def __init__(self, filename):
        self.filename = filename
        self.prs = Presentation()
        self.current_slide = None
        self.chart_type = XL_CHART_TYPE.BAR_CLUSTERED
        self.slide_top_offset = self.DEFAULT_TOP_OFFSET

    def add_slide(self):
    # Use a fully blank slide layout (no title box)
        self.current_slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide_top_offset = self.DEFAULT_TOP_OFFSET

    def ensure_space(self, element_height):
        if self.current_slide is None or (self.slide_top_offset + element_height > self.MAX_CONTENT_HEIGHT):
            self.add_slide()

    def add_title(self, text):
        height = Inches(0.5)
        self.ensure_space(height)

        left = Inches(0.5)
        width = Inches(9)
        top = self.slide_top_offset

        textbox = self.current_slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # vertical centering

        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER          # horizontal centering
        p.font.size = Pt(28)

        self.slide_top_offset += height + self.ELEMENT_SPACING

    def add_table(self, data):
        if not data:
            return

        headers = list(data[0].keys())
        rows = len(data) + 1
        cols = len(headers)
        row_height = 0.3
        table_height = Inches(0.5 + row_height * len(data))

        self.ensure_space(table_height)

        left = Inches(0.5)
        width = Inches(9)
        top = self.slide_top_offset

        table = self.current_slide.shapes.add_table(rows, cols, left, top, width, table_height).table

        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(18)

        for row_idx, row_data in enumerate(data, start=1):
            for col_idx, key in enumerate(headers):
                cell = table.cell(row_idx, col_idx)
                if isinstance(row_data[key], (int, float)):
                    cell.text = f"{row_data[key]:,.0f}"  # No decimal places
                else:
                    cell.text = str(row_data[key])
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER if col_idx == 1 else PP_ALIGN.LEFT

        self.slide_top_offset += table_height + Inches(1)  

    def add_graph(self, data):
        if not data:
            return

        headers = list(data[0].keys())
        label_key, value_key = headers[0], headers[1]

        categories = [d[label_key] for d in data]
        values = [d[value_key] for d in data]

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series(value_key, values)

        chart_height = Inches(4)
        chart_width = Inches(6)

        self.ensure_space(chart_height)

        x = Inches(1.5)
        y = self.slide_top_offset
        cx = chart_width
        cy = chart_height

        self.current_slide.shapes.add_chart(self.chart_type, x, y, cx, cy, chart_data)
        self.slide_top_offset += chart_height + self.ELEMENT_SPACING

    def jobs_cancelled_add_graph(self,data, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED):
        if not data:
            return

        if not self.current_slide:
            self.add_slide()

        #headers
        headers = list(data[0].keys())
        if len(headers) < 3:
            raise ValueError("Data must contain at least three keys per dictionary.")

        date_key, job_key, cancelled_key = headers[0], headers[1], headers[2]

        #data
        try:
            categories = [d[date_key] for d in data]
            job_total = [d[job_key] for d in data]
            job_cancelled = [d[cancelled_key] for d in data]
        except KeyError as e:
            raise ValueError(f"Key '{e.args[0]}' not found in data items.")

        #chart data
        chart_data = CategoryChartData()
        chart_data.categories = categories
        cancelled_total = sum(job_cancelled)
        job_total_sum = sum(job_total)

        #add series to chart data
        chart_data.add_series(cancelled_key, job_cancelled)
        chart_data.add_series(job_key, job_total)

        #chart positioning
        x, y, cx, cy = Inches(0.5), Inches(1), Inches(9), Inches(5.5)
        chart_frame = self.current_slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
        chart = chart_frame.chart

        left, top, width, height = Inches(6.5), Inches(6.2), Inches(3), Inches(1.2)


        def add_label_value_line(text_frame, label, value):
            p = text_frame.add_paragraph()
            p.alignment = PP_ALIGN.LEFT

            label_run = p.add_run()
            label_run.text = label
            label_run.font.size = Pt(8) if label in [
                '- Data shown includes  possible duplicate submission', 'Note:'
            ] else Pt(11)
            label_run.font.name = 'Arial'
            label_run.font.color.rgb = RGBColor(0, 0, 0)

            value_run = p.add_run()
            value_run.text = value
            value_run.font.size = Pt(13)
            value_run.font.name = 'Arial'
            value_run.font.color.rgb = RGBColor(0, 0, 0)
            value_run.font.bold = True

        #textbox
        textbox = self.current_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        # Title
        p_title = text_frame.paragraphs[0]
        p_title.clear()
        p_title.alignment = PP_ALIGN.CENTER

        title_run = p_title.add_run()
        title_run.text = "Total Job Received Count"
        title_run.font.size = Pt(14)
        title_run.font.name = 'Arial'
        title_run.font.color.rgb = RGBColor(0,0,0)
        title_run.font.bold = True

        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 255, 204)
        textbox.line.color.rgb = RGBColor(0, 0, 0)

        add_label_value_line(text_frame, "Total Jobs: ", f"{job_total_sum:,}")
        add_label_value_line(text_frame, "Jobs Done Within SLA: ", f"{cancelled_total:,}")
        add_label_value_line(text_frame, "Note:", "")
        add_label_value_line(text_frame, "- Data shown includes  possible duplicate submission", "")

        colors = [RGBColor(245, 105, 36), RGBColor(0, 66, 99)]

        #Category axis styling
        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(10)
        category_axis.tick_labels.font.name = 'Arial'

        #data labels styling
        for series in chart.series:
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.show_value = True
            data_labels.number_format = '#,##0'
            data_labels.font.size = Pt(12)
            data_labels.font.name = 'Arial'

        #series colors
        for idx, series in enumerate(chart.series):
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = colors[idx]


        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

                
    def add_SLA_graph(self, data, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED):
        if not data:
            return

        if not self.current_slide:
            self.add_slide()

        headers = list(data[0][0].keys())
        jobs = list(data[1][0].values())

        if len(headers) < 2:
            raise ValueError("Data must contain at least two keys per dictionary.")

        label_key, value_key = headers[0], headers[2]
        try:
            categories = [d[label_key] for d in data[0]]
            values = [d[value_key] for d in data[0]]
        except KeyError as e:
            raise ValueError(f"Key '{e.args[0]}' not found in data items.")

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series(value_key, values)

        x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(6)
        chart_frame = self.current_slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data)
        chart = chart_frame.chart
        #adds value at the top of graph
        for series in chart.series:
            series.has_data_labels = True
            data_labels = series.data_labels
            data_labels.show_value = True
            data_labels.number_format = '#,##0'
            data_labels.font.size = Pt(12)
            data_labels.font.name = 'Arial'

    def add_SLA_table(self, data):
        if not data:
            print("SLA table data is empty")
            return
        
        rows = len(data[0]) + 1
        cols = len(data[0][0])
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(0.8 + 0.3 * len(data[0]))

        table = self.current_slide.shapes.add_table(rows, cols, left, top, width, height).table

        #Header
        headers = list(data[0][0].keys())
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(18)

        #Rows
        for row_idx, row_data in enumerate(data[0], start=1):
            for col_idx, key in enumerate(headers):
                cell = table.cell(row_idx, col_idx)

                if isinstance(row_data[key], (int, float)):
                    cell.text = f"{row_data[key]:,.0f}"
                else:
                    cell.text = str(row_data[key])

                p = cell.text_frame.paragraphs[0]


                if col_idx == 1:
                    p.alignment = PP_ALIGN.CENTER
                else:
                    p.alignment = PP_ALIGN.LEFT  
                    
        summary_top = height + top + Inches(0.3)
        summary_left = left
        summary_width = Inches(4)
        summary_height = Inches(2)
        jobs = list(data[1][0].values())





        def add_label_value_line(text_frame, label, value):
            p = text_frame.add_paragraph()
            p.alignment = PP_ALIGN.LEFT

            #Label(normal)
            label_run = p.add_run()
            label_run.text = label
            if label != '(includes duplicate hash jobs)':
                label_run.font.size = Pt(14)
            else:
                label_run.font.size = Pt(10)
            label_run.font.name = 'Arial'
            label_run.font.color.rgb = RGBColor(0, 0, 0)


            #Value(bold)
            value_run = p.add_run()
            value_run.text = value
            value_run.font.size = Pt(14)
            value_run.font.name = 'Arial'
            value_run.font.color.rgb = RGBColor(0, 0, 0)
            value_run.font.bold = True

        #Setup the text box
        textbox = self.current_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            summary_left, summary_top, summary_width, summary_height
        )

        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        #First line centered
        p_title = text_frame.paragraphs[0]
        p_title.clear()
        p_title.alignment = PP_ALIGN.CENTER
        title_run = p_title.add_run()
        title_run.text = "Job SLA Status"
        title_run.font.size = Pt(16)
        title_run.font.name = 'Arial'
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(0, 0, 0)

        add_label_value_line(text_frame, "Total Done Jobs: ", f"{jobs[0]:,}")
        add_label_value_line(text_frame, "(includes duplicate hash jobs)", "")
        add_label_value_line(text_frame, "Jobs Done Within SLA: ", f"{jobs[1]:,}")
        add_label_value_line(text_frame, "Jobs Done Outside SLA: ", f"{jobs[0] - jobs[1]:,}")
        add_label_value_line(text_frame, "Overall SLA Compliance: ", f"{(jobs[1]/jobs[0])*100:.2f}%")

        note_left = summary_left + summary_width + Pt(20)
        note_top = summary_top
        note_width = Inches(4.5)
        note_height = summary_height

        note_textbox = self.current_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            note_left, note_top, note_width, note_height
        )

        note_text_frame = note_textbox.text_frame
        note_text_frame.clear()
        note_text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        note_fill = note_textbox.fill
        note_fill.solid()
        note_fill.fore_color.rgb = RGBColor(204, 255, 204)
        note_line = note_textbox.line
        note_line.color.rgb = RGBColor(0, 0, 0)

        #Note Title
        p_note_title = note_text_frame.paragraphs[0]
        p_note_title.clear()
        p_note_title.alignment = PP_ALIGN.LEFT
        note_run = p_note_title.add_run()
        note_run.text = "Note:"
        note_run.font.size = Pt(18)
        note_run.font.bold = True
        note_run.font.name = 'Arial'
        note_run.font.color.rgb = RGBColor(0, 0, 0)


        note_points = [
            "- Outside SLA are caused by delayed FRS scan results, need to manual run of Consumer Scanner & to manual restart of Metadata Extractor when it hangs (no assigned on weekends)",
            "- Auto-run & auto-restart have already been deployed in mid May 2025"
        ]

        for point in note_points:
            p = note_text_frame.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(12)
            p.font.name = 'Arial'
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.space_before = Pt(6)
            p.space_after = Pt(6)

        
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(204, 255, 204)
        line = textbox.line
        line.color.rgb = RGBColor(0, 0, 0)

    def save(self):
        self.prs.save(self.filename)
